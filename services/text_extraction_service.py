import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

class TextExtractionService:
    @staticmethod
    def extract_text_from_pdf(file_bytes):
        """
        Extracts raw text from the first few pages of a PDF.
        """
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            text = ""
            
            # Extract first 5 pages to avoid context bloat while getting enough info
            max_pages = min(len(reader.pages), 5)
            
            for i in range(max_pages):
                page_text = reader.pages[i].extract_text()
                if page_text:
                    text += page_text + "\n"
            
            return text.strip()
        except Exception as e:
            print(f"PDF Extraction Error: {e}")
            return None

    @staticmethod
    def extract_text_from_docx(file_bytes):
        """
        Extracts raw text from a Word document.
        """
        try:
            doc = Document(io.BytesIO(file_bytes))
            text = []
            for para in doc.paragraphs:
                if para.text:
                    text.append(para.text)
            
            # Limit to first 2000 words or so
            return "\n".join(text[:200]).strip() 
        except Exception as e:
            print(f"DOCX Extraction Error: {e}")
            return None

    @staticmethod
    def extract_text_from_pptx(file_bytes):
        """
        Extracts raw text from a PowerPoint presentation.
        """
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            text = []
            
            # Extract from first 10 slides
            max_slides = min(len(prs.slides), 10)
            
            for i in range(max_slides):
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            return "\n".join(text).strip()
        except Exception as e:
            print(f"PPTX Extraction Error: {e}")
            return None

    @staticmethod
    def extract_text(file_bytes, filename):
        """
        General extraction helper that routes by file extension.
        """
        ext = filename.split('.')[-1].lower()
        if ext == 'pdf':
            return TextExtractionService.extract_text_from_pdf(file_bytes)
        elif ext == 'docx':
            return TextExtractionService.extract_text_from_docx(file_bytes)
        elif ext == 'pptx':
            return TextExtractionService.extract_text_from_pptx(file_bytes)
        
        return None
